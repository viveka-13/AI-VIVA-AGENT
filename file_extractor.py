"""
file_extractor.py
-----------------
Utility to extract plain text from uploaded faculty materials.
Supports: PDF (.pdf), PowerPoint (.pptx), Word (.docx), Plain text (.txt)
"""
import os


def extract_text(filepath: str) -> str:
    """
    Extract all readable text from a file.
    Returns the extracted text as a single string.
    Raises ValueError for unsupported file types.
    """
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".pdf":
        return _extract_pdf(filepath)
    elif ext == ".pptx":
        return _extract_pptx(filepath)
    elif ext in (".docx",):
        return _extract_docx(filepath)
    elif ext == ".txt":
        return _extract_txt(filepath)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: .pdf, .pptx, .docx, .txt")


def _extract_pdf(filepath: str) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("PyMuPDF is required to read PDF files. Run: pip install PyMuPDF")

    text_parts = []
    with fitz.open(filepath) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts).strip()


def _extract_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required to read PPTX files. Run: pip install python-pptx")

    prs = Presentation(filepath)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
    return "\n".join(text_parts).strip()


def _extract_docx(filepath: str) -> str:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required to read DOCX files. Run: pip install python-docx")

    doc = Document(filepath)
    text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n".join(text_parts).strip()


def _extract_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()
